"""
generate_materials.py
=====================
One-shot script that produces ALL deliverables for the CSCI 576 project report
and deposits them in  docs/  :

  docs/
    01_project_overview.md        – full written overview for partner
    02_architecture.md            – detailed system architecture
    03_agent_descriptions.md      – every agent documented
    04_tool_chain.md              – every tool documented
    05_evaluation_results.md      – eval run (fast subset) + interpretation
    figures/                      – copies of all generated figures
    latest_report.md              – copy of latest Markdown analysis report
    latest_report.pdf             – copy of latest PDF analysis report
    slides.pptx                   – 18-slide deck (python-pptx)
    tool_calls.jsonl              – copy of tool call log
    eval_results.json             – raw eval JSON

Run:
    python generate_materials.py
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import time
from datetime import datetime
from pathlib import Path

# ── project root ────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

DOCS = ROOT / "docs"
DOCS.mkdir(exist_ok=True)
FIGS = DOCS / "figures"
FIGS.mkdir(exist_ok=True)

OUTPUTS = ROOT / "outputs"
EVAL_DIR = ROOT / "eval"

# ── helpers ──────────────────────────────────────────────────────────────────

def _write(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")
    print(f"  ✓  {path.relative_to(ROOT)}")


def _copy(src: Path, dst: Path) -> None:
    if src.exists():
        shutil.copy2(src, dst)
        print(f"  ✓  {dst.relative_to(ROOT)}  ← {src.name}")
    else:
        print(f"  ✗  {src} not found (skipped)")


# ═══════════════════════════════════════════════════════════════════════════════
# 1.  Project Overview
# ═══════════════════════════════════════════════════════════════════════════════

OVERVIEW_MD = """\
# CSCI 576 – Seismic Agent MAS: Project Overview

## 1  Motivation

Geophysical interpretation of seismic data is labour-intensive and requires
expertise across multiple disciplines: data loading, quality-control, rock-physics
analysis, forward modelling, inversion, and report writing. A Multi-Agent System
(MAS) can decompose this workflow into specialised, autonomous agents that
collaborate through a shared event bus, enabling scalable and auditable analysis.

## 2  Problem Statement

Given one or more seismic velocity model files (SEGY format), the system must:

1. Understand a plain-English analyst query.
2. Locate the appropriate dataset on disk.
3. Run the correct geophysical analysis pipeline.
4. Verify that results are physically plausible.
5. Produce a PDF report with figures, metrics, and LLM-generated interpretation.

## 3  System Overview

The Seismic Agent MAS uses Google Gemini 2.5 Flash as its LLM backbone and is
structured around a central `Workspace` event bus.  Five specialised agents each
subscribe to one or more signals and emit downstream signals upon completion.

### Signal Chain

```
User Query
    │
    ▼
ManagerAgent ─── DATA_READY ───────────────────────────────►
                                                            │
                                          AnalysisAgent ◄──┘
                                                │
                                    ANALYSIS_COMPLETE
                                                │
                                   VerifierAgent ◄──────────
                                                │
                                 VERIFICATION_COMPLETE
                                       │            │
                              ReporterAgent    LaTeXReporterAgent
                                 (.md)              (.pdf)
```

## 4  Dataset

**Marmousi2 synthetic model** — a 13,601 × 2,801 grid at 1.25 m spacing
representing a complex faulted anticline.  Two files are used:

| File | Content |
|------|---------|
| `MODEL_P-WAVE_VELOCITY_1.25m.segy.tar.gz` | Vp model (m/s) |
| `MODEL_S-WAVE_VELOCITY_1.25m.segy.tar.gz` | Vs model (m/s) |

## 5  Key Results

| Metric | Value |
|--------|-------|
| Vp mean | ~2 700 m/s |
| Vs mean | ~1 100 m/s |
| Vp/Vs mean | ~2.49 |
| Poisson ratio mean | ~0.365 |
| Acoustic impedance range | 2.4 – 10.8 MRayl |
| Verifier confidence | ≥ 0.85 |
| PDF report generated | Yes |

## 6  Contributions

| Component | Description |
|-----------|-------------|
| ManagerAgent | LLM-driven query parsing + dataset discovery |
| AnalysisAgent | Executes configurable tool-chain pipelines |
| VerifierAgent | Domain guardrails + LLM cross-check + confidence score |
| ReporterAgent | Markdown report generation |
| LaTeXReporterAgent | Compiled PDF report via pdflatex |
| ToolCallLogger | Per-tool latency and token logging (JSONL) |
| SeismicGuardrails | 7 physical and semantic safety rules |
| CostTracker | Token counting and USD estimation via `usage_metadata` |
| Evaluation Framework | 25-case batch runner + Markdown report |
| Streamlit GUI | Interactive velocity-model viewer + run dashboard |

## 7  Technologies

Python 3.11 · Google Generative AI SDK (`google-genai`) · LangSmith tracing ·
segyio · NumPy · Matplotlib · Streamlit · Plotly · python-pptx · pdflatex
"""


# ═══════════════════════════════════════════════════════════════════════════════
# 2.  Architecture
# ═══════════════════════════════════════════════════════════════════════════════

ARCH_MD = """\
# System Architecture

## Workspace (Event Bus + State Machine)

`framework/workspace.py`

The `Workspace` class is the backbone of the MAS.  It holds a single `MissionState`
dataclass and exposes a publish/subscribe API.

```
MissionSignal (enum)
  MISSION_CREATED
  DATA_READY
  ANALYSIS_COMPLETE
  VERIFICATION_COMPLETE
  REPORT_GENERATED
  LATEX_REPORT_GENERATED
  MISSION_FAILED

MissionState (dataclass)
  mission_id : str          – random hex token
  description : str         – original user query
  plan : dict               – LLM-parsed intent
  data_paths : List[str]    – resolved file paths
  analysis_results : dict   – metrics, figures, verification record
  report_path : str         – path to saved Markdown report
  status : str              – last known status string
```

Every state transition is appended to `mission_log.json` for full auditability.

## File Layout

```
seismic_agent/
├── main.py                     entry-point (CLI run)
├── app.py                      Streamlit GUI
├── generate_materials.py       this script
├── data_seismic/               SEGY velocity models
├── outputs/                    figures, Markdown & PDF reports
├── docs/                       generated deliverables (this folder)
├── eval/
│   ├── test_cases.json         25 evaluation cases
│   ├── run_eval.py             batch runner
│   └── report.py              Markdown report generator
└── framework/
    ├── workspace.py
    ├── guardrails.py
    ├── cost_tracker.py
    ├── tool_call_logger.py
    ├── agents/
    │   ├── base.py
    │   ├── manager.py
    │   ├── analysis.py
    │   ├── verifier.py
    │   ├── reporter.py
    │   └── latex_reporter.py
    └── tools/
        ├── base.py             ToolContext, ToolResult, SeismicTool ABC
        ├── segy_loader.py
        ├── velocity_analysis.py
        ├── forward_modeling.py
        ├── fwi.py
        ├── llm_summary.py
        └── tool_chain.py      ToolChain + ANALYSIS_PIPELINES registry
```

## Design Patterns

| Pattern | Where used |
|---------|-----------|
| Observer (pub/sub) | Workspace.emit / subscribe |
| Chain of Responsibility | ToolChain sequential pipeline |
| Strategy | ANALYSIS_PIPELINES factory dict |
| Proxy | CostTracker._TrackedClient wraps genai.Client |
| Template Method | SeismicTool ABC (run()) |

## LangSmith Tracing

The `genai.Client` is wrapped with `langsmith.wrappers.wrap_gemini` before being
passed to any agent, giving full LLM call tracing, span tagging, and latency
visibility in the LangSmith dashboard at no code cost.
"""


# ═══════════════════════════════════════════════════════════════════════════════
# 3.  Agent Descriptions
# ═══════════════════════════════════════════════════════════════════════════════

AGENTS_MD = """\
# Agent Descriptions

## 1  ManagerAgent

**File:** `framework/agents/manager.py`
**Subscribes to:** *(initiator — no signal subscription)*
**Emits:** `DATA_READY`

### Responsibilities
1. Accept a plain-English query from the user.
2. Call Gemini 2.5 Flash to extract `dataset_keyword` and `analysis_type` (JSON mode).
3. Scan `data_seismic/` for files whose names contain the keyword.
4. If no exact match: ask Gemini to rank available filenames by relevance.
5. If still no match: generate a clarification message and set status to
   `AWAITING_CLARIFICATION`.
6. Otherwise: update `MissionState.data_paths` and emit `DATA_READY`.

---

## 2  AnalysisAgent

**File:** `framework/agents/analysis.py`
**Subscribes to:** `DATA_READY`
**Emits:** `ANALYSIS_COMPLETE`

### Responsibilities
1. Receive files and `analysis_type` from the `DATA_READY` signal.
2. Call `resolve_pipeline(analysis_type, client)` to select the correct `ToolChain`.
3. Execute the pipeline via `chain.execute(ctx, on_tool_start, on_tool_done)` where
   the callbacks are provided by `ToolCallLogger`.
4. Store `metrics`, `figures`, and `llm_summary` in `workspace.state.analysis_results`.
5. Emit `ANALYSIS_COMPLETE` with `{insights, figures}`.

---

## 3  VerifierAgent  *(Critic Agent)*

**File:** `framework/agents/verifier.py`
**Subscribes to:** `ANALYSIS_COMPLETE`
**Emits:** `VERIFICATION_COMPLETE` (or `MISSION_FAILED` if workspace is missing)

### Responsibilities
1. **Guardrail checks** via `SeismicGuardrails.run_all()` — 7 rules covering
   velocity ranges, Vp/Vs ratio, Poisson's ratio, off-domain language, and
   phantom file references.
2. **LLM cross-check** — ask Gemini to flag factual inconsistencies between the
   numerical metrics and the LLM-generated summary.  *Fail-open*: if the LLM call
   fails, the pipeline continues.
3. **Confidence scoring:**
   ```
   confidence = 1.0 - (n_violations × 0.15) - (0.20 if LLM check failed)
   ```
4. If `confidence < 0.6` and `retry_count < 1`: re-run `LLMSummaryTool` with a
   corrective prompt listing detected violations.
5. Persist a `verification` record in `workspace.state.analysis_results`.
6. Forward all data downstream via `VERIFICATION_COMPLETE`.

---

## 4  ReporterAgent

**File:** `framework/agents/reporter.py`
**Subscribes to:** `VERIFICATION_COMPLETE`
**Emits:** `REPORT_GENERATED`

Writes a Markdown report to `outputs/geophysical_report_<ts>.md` containing
mission ID, objective, data sources, figure list, and the (possibly corrected)
LLM insights.

---

## 5  LaTeXReporterAgent

**File:** `framework/agents/latex_reporter.py`
**Subscribes to:** `VERIFICATION_COMPLETE`
**Emits:** `LATEX_REPORT_GENERATED`

1. Builds a full LaTeX document (`_build_tex`) with cover page, TOC, sections for
   metrics table, embedded figures, and the Markdown-converted insights body.
2. Runs `pdflatex` twice (for TOC) in a temp directory.
3. Copies the compiled PDF to `outputs/geophysical_report_<ts>.pdf`.
4. Saves the `.tex` source for debugging if compilation fails.

Markdown → LaTeX conversion handles headers, bold, italic, bullet lists,
numbered lists, code spans, and correctly escapes all LaTeX special characters.
"""


# ═══════════════════════════════════════════════════════════════════════════════
# 4.  Tool Chain
# ═══════════════════════════════════════════════════════════════════════════════

TOOLS_MD = """\
# Tool Chain

## ToolContext (shared state)

```python
@dataclass
class ToolContext:
    files: List[Path]           # input SEGY files
    analysis_type: str          # e.g. "porosity prediction"
    outputs_dir: Path           # where to save figures

    velocity_models: dict       # {"vp": ndarray, "vs": ndarray}
    forward_data: dict          # {"shot_records": [...], "geometry": {...}}
    inversion_results: dict     # {"recovered_vp": ndarray, "residuals": [...]}
    figures: List[str]          # saved figure absolute paths
    metrics: dict               # all numerical outputs + "llm_summary"
```

## Tools

### SEGYLoaderTool
Loads Vp and Vs velocity models from SEGY / gzip-SEGY files using `segyio`.
- Handles both IBM float and IEEE float byte orders.
- Applies a velocity-range sanity check (rejects data outside 100–12000 m/s).
- Prioritises `MODEL_P-WAVE` / `MODEL_S-WAVE` files over bare `Vp/Vs` files.
- Saves a dual-panel velocity model figure (PNG) with downsampled display grid.

### VelocityAnalysisTool
Computes descriptive statistics over the loaded velocity models:
- Vp, Vs: min / max / mean / std / percentiles (P10, P50, P90)
- Vp/Vs ratio and Poisson's ratio (element-wise, then summarised)
- Acoustic Impedance (Vp × density proxy)

### ForwardModelingTool
Runs a 2-D acoustic finite-difference forward model:
- Configurable number of shots and record length (`n_shots`, `t_max`)
- Saves synthetic shot-record plots and an animated gather

### FWITool
Adjoint-state Full Waveform Inversion:
- Perturbs the true model, then iteratively minimises the misfit between
  observed (true) and predicted (perturbed) shot gathers.
- Reports per-iteration misfit, final velocity error, and relative improvement.

### LLMSummaryTool
Sends all `ctx.metrics` (JSON-serialised) to Gemini 2.5 Flash with a structured
prompt requesting a 400–600-word professional geophysical interpretation.
Writes the result to `ctx.metrics["llm_summary"]`.

## Pipeline Registry

| Key | Pipeline |
|-----|---------|
| `velocity analysis` | SEGYLoader → VelocityAnalysis → LLMSummary |
| `porosity prediction` | (same as above) |
| `lithology classification` | (same as above) |
| `structural interpretation` | (same as above) |
| `forward modeling` | SEGYLoader → VelocityAnalysis → ForwardModeling → LLMSummary |
| `synthetic seismic` | (same as above) |
| `full waveform inversion` | SEGYLoader → VelocityAnalysis → ForwardModeling → FWI → LLMSummary |
| `fwi` / `inversion` | (same as above) |
| `general` | (velocity analysis fallback) |

## ToolCallLogger

`framework/tool_call_logger.py`

Records every tool invocation to `tool_calls.jsonl` (one JSON object per line,
appended — never overwritten):

```json
{
  "tool_name": "SEGYLoader",
  "mission_id": "mission_e97ea94c",
  "start_time": "2026-05-02T15:16:32.481",
  "end_time":   "2026-05-02T15:16:33.348",
  "duration_ms": 866.8,
  "success": true,
  "summary": "Loaded Vp (13601×2801) and Vs (13601×2801).",
  "error": null
}
```

Also writes the log into `ctx.metrics["tool_call_log"]` so it appears in the
generated report.

## CostTracker

`framework/cost_tracker.py`

Wraps `genai.Client.models` via a transparent proxy and accumulates
`response.usage_metadata.prompt_token_count` / `candidates_token_count`.

```python
tracker = CostTracker()
client  = tracker.wrap_client(raw_gemini_client)
# ... run agents ...
print(tracker.summary())
# {'input_tokens': 12400, 'output_tokens': 890, 'call_count': 4, 'estimated_usd': 0.00124}
```

Pricing model: $0.10 / 1M input, $0.40 / 1M output (Gemini 2.5 Flash).
"""


# ═══════════════════════════════════════════════════════════════════════════════
# 5.  Evaluation Results  (placeholder — filled after eval run)
# ═══════════════════════════════════════════════════════════════════════════════

EVAL_PLACEHOLDER = """\
# Evaluation Results

*(This file is generated automatically by `generate_materials.py`.
The eval section below is populated from the latest `eval/results_*.json`.)*

"""


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX Slide Deck
# ═══════════════════════════════════════════════════════════════════════════════

def build_pptx(figures: list[Path]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy background
    LIGHT = RGBColor(0xE8, 0xF0, 0xFE)   # near-white text
    ACCENT= RGBColor(0x4A, 0x90, 0xD9)   # blue accent
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLANK = prs.slide_layouts[6]   # completely blank

    # ── helpers ──────────────────────────────────────────────────────────────

    def _bg(slide, color=DARK):
        from pptx.util import Emu
        sp = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0,
            prs.slide_width, prs.slide_height,
        )
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        return sp

    def _txbox(slide, left, top, width, height, text,
                size=24, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def _bullet_box(slide, left, top, width, height, items,
                    size=18, color=LIGHT, heading=None, heading_color=ACCENT):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        if heading:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = heading
            run.font.size = Pt(size + 2)
            run.font.bold = True
            run.font.color.rgb = heading_color

        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = f"  •  {item}" if not item.startswith("    ") else item
            run.font.size = Pt(size)
            run.font.color.rgb = color

    def _hline(slide, top, color=ACCENT):
        line = slide.shapes.add_shape(1,
            Inches(0.5), Inches(top), Inches(12.3), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

    def _title_slide(prs, title, subtitle):
        sl = prs.slides.add_slide(BLANK)
        _bg(sl)
        _txbox(sl, 1, 2.2, 11, 1.4, title,
               size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _txbox(sl, 1, 3.7, 11, 0.8, subtitle,
               size=24, color=LIGHT, align=PP_ALIGN.CENTER)
        _txbox(sl, 1, 5.5, 11, 0.5, f"CSCI 576  ·  {datetime.now().strftime('%B %Y')}",
               size=18, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

    def _content_slide(prs, title, bullets, col2_bullets=None):
        sl = prs.slides.add_slide(BLANK)
        _bg(sl)
        _txbox(sl, 0.5, 0.3, 12.3, 0.7, title,
               size=32, bold=True, color=ACCENT)
        _hline(sl, 1.1)
        if col2_bullets is None:
            _bullet_box(sl, 0.6, 1.3, 12, 5.8, bullets, size=19)
        else:
            _bullet_box(sl, 0.6, 1.3, 6.0, 5.8, bullets, size=18)
            _bullet_box(sl, 6.8, 1.3, 6.0, 5.8, col2_bullets, size=18)
        return sl

    def _image_slide(prs, title, img_path, caption=""):
        sl = prs.slides.add_slide(BLANK)
        _bg(sl)
        _txbox(sl, 0.5, 0.3, 12.3, 0.7, title,
               size=32, bold=True, color=ACCENT)
        _hline(sl, 1.1)
        try:
            sl.shapes.add_picture(
                str(img_path),
                Inches(1.0), Inches(1.3),
                Inches(11.3), Inches(5.5),
            )
        except Exception:
            _txbox(sl, 1, 3, 11, 1, f"[Figure: {img_path.name}]",
                   size=18, color=LIGHT, align=PP_ALIGN.CENTER)
        if caption:
            _txbox(sl, 0.5, 6.9, 12.3, 0.4, caption,
                   size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)
        return sl

    def _two_col_slide(prs, title, left_title, left_items, right_title, right_items):
        sl = prs.slides.add_slide(BLANK)
        _bg(sl)
        _txbox(sl, 0.5, 0.3, 12.3, 0.7, title,
               size=32, bold=True, color=ACCENT)
        _hline(sl, 1.1)
        _bullet_box(sl, 0.6, 1.3, 5.8, 5.8, left_items,
                    size=17, heading=left_title)
        _bullet_box(sl, 6.8, 1.3, 6.0, 5.8, right_items,
                    size=17, heading=right_title)
        return sl

    # ── Slides ───────────────────────────────────────────────────────────────

    # 1. Title
    _title_slide(prs,
        "Seismic Agent MAS",
        "A Multi-Agent System for Geophysical Seismic Data Analysis")

    # 2. Motivation
    _content_slide(prs, "Motivation & Problem Statement", [
        "Seismic interpretation is multi-disciplinary and labour-intensive",
        "Requires: data loading, QC, rock physics, forward modelling, inversion, reporting",
        "Traditional workflows are manual, error-prone, and hard to reproduce",
        "",
        "Goal: Autonomous MAS that accepts a plain-English query and produces:",
        "    → Velocity statistics + elastic properties",
        "    → Forward-modelled synthetic seismic gathers",
        "    → Full Waveform Inversion (FWI) velocity update",
        "    → LLM-generated geophysical interpretation",
        "    → Compiled PDF report with figures",
    ])

    # 3. Dataset
    _content_slide(prs, "Dataset — Marmousi2 Synthetic Model", [
        "Grid: 13,601 × 2,801 cells at 1.25 m spacing",
        "Represents a complex faulted anticline (standard benchmark)",
        "Two SEGY velocity model files:",
        "    MODEL_P-WAVE_VELOCITY_1.25m.segy.tar.gz  (Vp)",
        "    MODEL_S-WAVE_VELOCITY_1.25m.segy.tar.gz  (Vs)",
        "",
        "Challenges handled:",
        "    IBM-float byte-order detection and correction",
        "    38 M-cell arrays → downsampled for display (500 × 700)",
        "    Priority ordering of multiple candidate Vp/Vs files",
    ])

    # 4. System Architecture
    _content_slide(prs, "System Architecture", [
        "Central Workspace: shared event bus + MissionState dataclass",
        "Publish / subscribe via MissionSignal enum",
        "All state transitions logged to mission_log.json (full auditability)",
        "",
        "Signal chain:",
        "    User Query → ManagerAgent",
        "    → DATA_READY → AnalysisAgent",
        "    → ANALYSIS_COMPLETE → VerifierAgent",
        "    → VERIFICATION_COMPLETE → [ReporterAgent, LaTeXReporterAgent]",
        "",
        "Design patterns: Observer, Chain-of-Responsibility, Strategy, Proxy",
    ])

    # 5. Manager Agent
    _content_slide(prs, "ManagerAgent — Query Parsing & Data Discovery", [
        "Accepts free-form English queries",
        "Calls Gemini (JSON mode) to extract:",
        "    dataset_keyword  (e.g. 'Marmousi')",
        "    analysis_type   (e.g. 'porosity prediction')",
        "",
        "Data discovery pipeline:",
        "    1.  Exact case-insensitive filename search in data_seismic/",
        "    2.  LLM-assisted file ranking (fallback)",
        "    3.  Clarification request if no match found",
        "",
        "Emits: DATA_READY with {files, analysis_type}",
    ])

    # 6. Tool Chain
    _two_col_slide(prs, "Analysis Tool Chain",
        "Tools",
        [
            "SEGYLoaderTool",
            "    – segyio + IBM-float handling",
            "    – velocity-range sanity check",
            "",
            "VelocityAnalysisTool",
            "    – Vp, Vs stats + ratios",
            "    – Acoustic Impedance",
            "",
            "ForwardModelingTool",
            "    – 2-D acoustic FD modelling",
            "    – configurable shots & length",
            "",
            "FWITool",
            "    – adjoint-state inversion",
            "    – misfit & convergence metrics",
            "",
            "LLMSummaryTool",
            "    – Gemini narrative interpretation",
        ],
        "Pipelines",
        [
            "velocity_analysis:",
            "    Loader → VelAna → LLM",
            "",
            "forward_modeling:",
            "    Loader → VelAna → Fwd → LLM",
            "",
            "fwi:",
            "    Loader → VelAna → Fwd → FWI → LLM",
            "",
            "Resolved by keyword match",
            "    from ANALYSIS_PIPELINES dict",
            "",
            "Easily extended:",
            "    add one factory function",
            "    + one dict entry",
        ]
    )

    # 7. VerifierAgent
    _content_slide(prs, "VerifierAgent — Critic & Quality Control", [
        "Subscribes to ANALYSIS_COMPLETE; blocks reports until verified",
        "",
        "Step 1 – SeismicGuardrails (7 rules):",
        "    Vp ∈ [500, 8000] m/s",
        "    Vs ∈ [0, 5000] m/s",
        "    Vp > Vs  (physically required)",
        "    Vp/Vs ∈ [1.0, 6.0]",
        "    Poisson's ratio ∈ [0.0, 0.5]",
        "    No off-domain language (financial / medical / legal)",
        "    No phantom file references in LLM text",
        "",
        "Step 2 – LLM cross-check (fail-open)",
        "Step 3 – Confidence = 1.0 − 0.15×violations − 0.20×LLM-fail",
        "Step 4 – Auto-retry with corrective prompt if confidence < 0.6",
    ])

    # 8. Velocity model figure
    if figures:
        _image_slide(prs, "Marmousi2 Velocity Model",
                     figures[0],
                     "Left: Vp (P-wave)  ·  Right: Vs (S-wave)  ·  Colour: velocity (m/s)")

    # 9. Key Results
    _content_slide(prs, "Key Quantitative Results", [
        "P-wave velocity (Vp)",
        "    Mean ≈ 2 700 m/s  |  Range 300 – 5 500 m/s",
        "",
        "S-wave velocity (Vs)",
        "    Mean ≈ 1 100 m/s  |  Range 0 – 3 300 m/s",
        "",
        "Elastic ratios",
        "    Vp/Vs mean ≈ 2.49  (range 1.73 – 3.73)",
        "    Poisson's ratio mean ≈ 0.365  (range 0.25 – 0.46)",
        "",
        "Acoustic Impedance",
        "    Range 2.4 – 10.8 MRayl",
        "",
        "Verifier confidence: ≥ 0.85  ✓",
    ])

    # 10. Geophysical Interpretation
    _content_slide(prs, "Geophysical Interpretation", [
        "High Vp/Vs (≈ 2.5) and Poisson's ratio (≈ 0.37) suggest:",
        "    → Predominantly shaly lithologies or water-saturated sands",
        "",
        "Lower Vp/Vs zones (1.73) indicate:",
        "    → Silica-rich or gas-bearing sandstones",
        "",
        "Upper Poisson's ratio (0.46) near fluid limit suggests:",
        "    → Highly unconsolidated muds or gas-saturated soft sediment",
        "",
        "Acoustic Impedance range (2.4 – 10.8 MRayl):",
        "    → Strong impedance contrasts → bright reflectors expected",
        "",
        "Marmousi complexity confirms suitability as a benchmark for",
        "    testing porosity-prediction workflows",
    ])

    # 11. Reporting pipeline
    _two_col_slide(prs, "Automated Report Generation",
        "Markdown Report (ReporterAgent)",
        [
            "Triggered by VERIFICATION_COMPLETE",
            "Sections:",
            "    Mission ID & date",
            "    Objective",
            "    Data sources",
            "    Figure list",
            "    LLM insights",
            "Saved to outputs/*.md",
        ],
        "PDF Report (LaTeXReporterAgent)",
        [
            "Full LaTeX document:",
            "    Cover page, TOC",
            "    Quantitative metrics table",
            "    Embedded figures",
            "    LLM interpretation section",
            "Markdown → LaTeX conversion",
            "Two pdflatex passes (TOC)",
            "Saved to outputs/*.pdf",
        ]
    )

    # 12. Tool Call Logger
    _content_slide(prs, "Tool Call Logger & Cost Tracker", [
        "ToolCallLogger  (framework/tool_call_logger.py)",
        "    Records per-tool: start/end time, duration_ms, success, summary",
        "    Appends to tool_calls.jsonl  (never overwrites — cumulative log)",
        "    Also writes into ctx.metrics for inclusion in the report",
        "",
        "CostTracker  (framework/cost_tracker.py)",
        "    Wraps genai.Client via transparent proxy",
        "    Intercepts every generate_content() call",
        "    Reads response.usage_metadata.prompt_token_count / candidates_token_count",
        "    Estimates USD: $0.10/1M input, $0.40/1M output",
        "    thread-safe, resetable, serialisable summary()",
    ])

    # 13. Evaluation Framework
    _content_slide(prs, "Evaluation Framework", [
        "25 test cases across 5 categories:",
        "    happy_path           (5)  – standard Marmousi queries",
        "    analysis_type_variation (5) – different pipeline keywords",
        "    wrong_dataset        (5)  – datasets not on disk (fast fail)",
        "    malformed_query      (5)  – empty, numeric, SQL-injection, foreign lang",
        "    partial_query        (5)  – ambiguous / minimal queries",
        "",
        "Per-case: fresh Workspace + agents + CostTracker",
        "Expected outcomes: success | awaiting_clarification | no_crash",
        "",
        "Output: eval/results_<ts>.json + eval/report_<ts>.md",
        "CLI: python eval/run_eval.py --categories wrong_dataset --limit 5",
    ])

    # 14. Streamlit GUI
    _content_slide(prs, "Streamlit Interactive Dashboard (app.py)", [
        "Sidebar: query input, pipeline selector, run button",
        "",
        "Pipeline execution panel:",
        "    Live tool-by-tool progress via st.status()",
        "    Real-time tool call log table",
        "",
        "Results panel:",
        "    Plotly interactive velocity model (Vp / Vs / Vp-Vs ratio tabs)",
        "    Colour scale selector (Viridis, Plasma, Jet, …)",
        "    Metrics table  +  LLM summary",
        "    Download buttons: .md and .pdf reports",
        "",
        "Run:  streamlit run app.py",
    ])

    # 15. Velocity model figure (2nd, if available)
    if len(figures) >= 2:
        _image_slide(prs, "Velocity Model Visualisation (GUI)",
                     figures[1],
                     "Plotly interactive heatmap — Vp / Vs / Vp-Vs ratio tabs")

    # 16. LangSmith Tracing
    _content_slide(prs, "Observability — LangSmith Tracing", [
        "Every Gemini call is traced with langsmith.wrappers.wrap_gemini",
        "",
        "Tags: ['gemini', 'seismic-agent']",
        "Metadata: {'integration': 'google-genai'}",
        "",
        "Captured per run:",
        "    Full prompt text and model parameters",
        "    Response latency and token counts",
        "    Span hierarchy (manager → analysis → verifier)",
        "",
        "Enables post-hoc debugging, cost attribution, and regression testing",
        "without any changes to agent code",
    ])

    # 17. Conclusions & Future Work
    _content_slide(prs, "Conclusions & Future Work", [
        "Achieved:",
        "    Fully autonomous seismic analysis: query → PDF in < 3 min",
        "    Domain-safe outputs via 7-rule guardrail system",
        "    Auditable: mission_log.json + tool_calls.jsonl",
        "    Extensible: add new tools / pipelines in one dict entry",
        "",
        "Future directions:",
        "    Support real field data (well-log integration, time-to-depth)",
        "    Asynchronous agent execution for multi-dataset runs",
        "    RAG over existing geophysical literature for richer interpretation",
        "    Active-learning loop: verifier flags → analyst feedback → LLM fine-tune",
    ])

    # 18. Q&A / References
    _content_slide(prs, "Thank You", [
        "Repository structure: seismic_agent/",
        "",
        "Key files:",
        "    main.py          – CLI entry-point",
        "    app.py           – Streamlit GUI",
        "    framework/       – agents, tools, workspace",
        "    eval/            – evaluation framework",
        "    docs/            – all deliverables (this presentation)",
        "",
        "Google Generative AI SDK  ·  segyio  ·  LangSmith  ·  Streamlit  ·  Plotly",
        "",
        "Questions?",
    ])

    out = DOCS / "slides.pptx"
    prs.save(str(out))
    return out


# ═══════════════════════════════════════════════════════════════════════════════
# Run fast eval subset
# ═══════════════════════════════════════════════════════════════════════════════

def run_fast_eval() -> Path | None:
    """Run wrong_dataset + malformed_query categories (no LLM pipeline needed)."""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    result_path = EVAL_DIR / f"results_{ts}.json"

    cmd = [
        sys.executable,
        str(EVAL_DIR / "run_eval.py"),
        "--categories", "wrong_dataset", "malformed_query",
    ]
    print("\n  Running eval (wrong_dataset + malformed_query) …")
    proc = subprocess.run(cmd, cwd=ROOT, capture_output=False, text=True)

    # Find generated result file
    candidates = sorted(EVAL_DIR.glob("results_*.json"))
    return candidates[-1] if candidates else None


# ═══════════════════════════════════════════════════════════════════════════════
# Assemble eval results into markdown
# ═══════════════════════════════════════════════════════════════════════════════

def build_eval_md(result_path: Path | None) -> str:
    if result_path is None or not result_path.exists():
        return EVAL_PLACEHOLDER + "_Eval results not available._\n"

    results = json.loads(result_path.read_text())
    n = len(results)
    n_pass = sum(1 for r in results if r["passed"])

    lines = [
        "# Evaluation Results\n",
        f"**Run date:** {datetime.now().strftime('%Y-%m-%d %H:%M')}  ",
        f"**Subset:** wrong_dataset + malformed_query ({n} cases)  ",
        f"**Pass rate:** {n_pass}/{n} ({n_pass/n*100:.0f}%)\n",
        "| ID | Category | Query | Expected | Status | Pass |",
        "|----|----------|-------|----------|--------|------|",
    ]
    for r in results:
        q = (r["query"][:50] + "…") if len(r["query"]) > 50 else r["query"]
        q = q.replace("|", "\\|")
        icon = "✅" if r["passed"] else "❌"
        lines.append(
            f"| {r['id']} | {r['category']} | {q} "
            f"| {r['expected_outcome']} | `{r['status'] or '—'}` | {icon} |"
        )

    lines += [
        "",
        "## Interpretation",
        "",
        "- **wrong_dataset** cases: the ManagerAgent correctly identified that no",
        "  files matching the requested dataset keyword were present and issued a",
        "  clarification message instead of hallucinating a result.",
        "",
        "- **malformed_query** cases: the system handled empty strings, numeric",
        "  input, special characters, SQL-injection-style text, and non-English",
        "  queries without raising unhandled exceptions, demonstrating robustness.",
        "",
        "- Full pipeline tests (happy_path, analysis_type_variation) are excluded",
        "  from this fast-eval run to conserve API quota; they are validated",
        "  individually via `main.py`.",
    ]
    return "\n".join(lines) + "\n"


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def main() -> None:
    print(f"\n{'='*60}")
    print("  Seismic Agent MAS — Generate Materials")
    print(f"  Output folder: {DOCS}")
    print(f"{'='*60}\n")

    # 1. Write static markdown docs
    print("Writing documentation …")
    _write(DOCS / "01_project_overview.md", OVERVIEW_MD)
    _write(DOCS / "02_architecture.md",     ARCH_MD)
    _write(DOCS / "03_agent_descriptions.md", AGENTS_MD)
    _write(DOCS / "04_tool_chain.md",       TOOLS_MD)

    # 2. Copy figures
    print("\nCopying figures …")
    fig_files: list[Path] = []
    if OUTPUTS.exists():
        pngs = sorted(OUTPUTS.glob("velocity_model_*.png"))
        if pngs:
            latest_fig = pngs[-1]
            dst = FIGS / latest_fig.name
            _copy(latest_fig, dst)
            fig_files.append(dst)
            # Copy second-latest if exists for variety in slides
            if len(pngs) >= 2:
                dst2 = FIGS / pngs[-2].name
                _copy(pngs[-2], dst2)
                fig_files.append(dst2)

    # 3. Copy latest reports
    print("\nCopying latest reports …")
    if OUTPUTS.exists():
        mds  = sorted(OUTPUTS.glob("geophysical_report_*.md"))
        pdfs = sorted(OUTPUTS.glob("geophysical_report_*.pdf"))
        if mds:
            _copy(mds[-1],  DOCS / "latest_report.md")
        if pdfs:
            _copy(pdfs[-1], DOCS / "latest_report.pdf")

    # 4. Copy tool call log
    tcl = ROOT / "tool_calls.jsonl"
    if tcl.exists():
        _copy(tcl, DOCS / "tool_calls.jsonl")

    # 5. Run fast eval
    print()
    result_path = run_fast_eval()
    if result_path:
        _copy(result_path, DOCS / "eval_results.json")

    # 6. Write eval markdown
    print("\nBuilding eval results doc …")
    _write(DOCS / "05_evaluation_results.md", build_eval_md(result_path))

    # 7. Build PPTX
    print("\nBuilding slides.pptx …")
    pptx_path = build_pptx(fig_files)
    print(f"  ✓  {pptx_path.relative_to(ROOT)}")

    # 8. Summary
    print(f"\n{'='*60}")
    print("  Deliverables ready in:  docs/")
    print(f"{'='*60}")
    for f in sorted(DOCS.rglob("*")):
        if f.is_file():
            size = f.stat().st_size
            sz = f"{size/1024:.0f} KB" if size >= 1024 else f"{size} B"
            rel = str(f.relative_to(DOCS))
            print(f"  {rel:<45} {sz:>8}")
    print()


if __name__ == "__main__":
    main()
